from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== Color Palette =====
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
LIGHT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
ACCENT_GREEN = RGBColor(0x52, 0xC4, 0x1A)
ACCENT_ORANGE = RGBColor(0xFA, 0x8C, 0x16)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)

output_dir = r"C:\Users\27125\Desktop\lzy的Javaweb作业\library"


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=DARK_GRAY, line_spacing=1.5):
    """lines is a list of (text, bold, color_override, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, clr, fs = line_data, False, color, font_size
        else:
            text = line_data[0]
            bold = line_data[1] if len(line_data) > 1 else False
            clr = line_data[2] if len(line_data) > 2 else color
            fs = line_data[3] if len(line_data) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


# ============================================================
# Slide 1: Cover
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BLUE)
# Decorative bars
add_rect(slide1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), LIGHT_BLUE)
add_rect(slide1, Inches(0), Inches(4.35), Inches(13.333), Inches(0.04), RGBColor(0x42, 0xA5, 0xF5))
# Title
add_text_box(slide1, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
             "图书馆借还系统", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# Subtitle
add_text_box(slide1, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.7),
             "基于 Java Servlet + JSP + MySQL 的图书借阅归还管理平台", font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
# Info
add_multi_text(slide1, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.8), [
    ("JavaWeb 期末答辩", True, WHITE, 24),
    ("", False, WHITE, 8),
    ("专    业：计算机科学与技术", False, RGBColor(0xE3, 0xF2, 0xFD), 18),
    ("指导教师：_______________", False, RGBColor(0xE3, 0xF2, 0xFD), 18),
    ("答辩日期：2026年6月", False, RGBColor(0xE3, 0xF2, 0xFD), 18),
], font_size=18, color=RGBColor(0xE3, 0xF2, 0xFD))

# ============================================================
# Slide 2: Table of Contents
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)
add_rect(slide2, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide2, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide2, Inches(1.3), Inches(1.5), Inches(5), Inches(0.7),
             "目  录  CONTENTS", font_size=32, color=DARK_BLUE, bold=True)

items = [
    ("01", "项目背景与意义"),
    ("02", "技术栈介绍"),
    ("03", "系统架构设计"),
    ("04", "数据库设计"),
    ("05", "功能模块展示"),
    ("06", "核心代码解析"),
    ("07", "项目总结与展望"),
]
for i, (num, title) in enumerate(items):
    y = 2.6 + i * 0.65
    add_rect(slide2, Inches(1.6), Inches(y), Inches(0.5), Inches(0.5), LIGHT_BLUE if i == 0 else LIGHT_GRAY)
    add_text_box(slide2, Inches(1.6), Inches(y + 0.05), Inches(0.5), Inches(0.4),
                 num, font_size=18, color=WHITE if i == 0 else DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide2, Inches(2.4), Inches(y + 0.05), Inches(8), Inches(0.4),
                 title, font_size=18, color=DARK_GRAY)

# ============================================================
# Slide 3: Project Background
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)
add_rect(slide3, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide3, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide3, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "01  项目背景与意义", font_size=32, color=DARK_BLUE, bold=True)

# Left column - Background
add_text_box(slide3, Inches(1.3), Inches(2.5), Inches(5.5), Inches(0.5),
             "📋 项目背景", font_size=22, color=LIGHT_BLUE, bold=True)
add_multi_text(slide3, Inches(1.3), Inches(3.1), Inches(5.5), Inches(4.0), [
    ("• 高校图书馆藏书量日益增长，传统人工管理效率低下", False, DARK_GRAY, 16),
    ("• 借还书记录容易遗漏、难以追溯", False, DARK_GRAY, 16),
    ("• 读者无法实时查询图书库存状态", False, DARK_GRAY, 16),
    ("• 需要一个轻量级的数字化管理工具", False, DARK_GRAY, 16),
], font_size=16)

# Right column - Significance
add_text_box(slide3, Inches(7.3), Inches(2.5), Inches(5.5), Inches(0.5),
             "🎯 项目意义", font_size=22, color=ACCENT_GREEN, bold=True)
add_multi_text(slide3, Inches(7.3), Inches(3.1), Inches(5.5), Inches(4.0), [
    ("✓ 实现在线图书检索与借阅", False, DARK_GRAY, 16),
    ("✓ 自动记录借还信息，告别纸质登记", False, DARK_GRAY, 16),
    ("✓ 管理员可高效管理图书和借阅记录", False, DARK_GRAY, 16),
    ("✓ 读者可随时查看自己的借阅历史", False, DARK_GRAY, 16),
    ("✓ 系统轻量，部署方便，适合小型图书馆", False, DARK_GRAY, 16),
], font_size=16)

# ============================================================
# Slide 4: Tech Stack
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)
add_rect(slide4, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide4, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide4, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "02  技术栈介绍", font_size=32, color=DARK_BLUE, bold=True)

techs = [
    ("☕ Java 8+", "核心开发语言\nServlet + JSP + JDBC", DARK_BLUE),
    ("🗄️ MySQL 8.0", "关系型数据库\nInnoDB引擎 + 外键约束", RGBColor(0xE0, 0x7A, 0x00)),
    ("🚀 Jetty 9.4", "嵌入式Web服务器\n零配置，即跑即用", ACCENT_GREEN),
    ("🎨 HTML5+CSS3", "前端页面构建\nFlexbox + Grid 响应式", RGBColor(0xE3, 0x4E, 0x26)),
    ("📦 Maven", "项目构建与依赖管理\n一键编译运行", RGBColor(0xC7, 0x1A, 0x36)),
    ("🔄 JSTL + EL", "JSP标签库\n简化动态页面开发", RGBColor(0x6A, 0x1B, 0x9A)),
]
for i, (title, desc, color) in enumerate(techs):
    x = 1.2 + (i % 3) * 3.9
    y = 2.6 + (i // 3) * 2.3
    # Card bg
    card = add_rect(slide4, Inches(x), Inches(y), Inches(3.4), Inches(1.9), LIGHT_GRAY)
    # Accent bar
    add_rect(slide4, Inches(x), Inches(y), Inches(3.4), Inches(0.06), color)
    add_text_box(slide4, Inches(x + 0.3), Inches(y + 0.3), Inches(3.0), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_multi_text(slide4, Inches(x + 0.3), Inches(y + 0.9), Inches(3.0), Inches(1.0), [
        (desc, False, MEDIUM_GRAY, 14),
    ])

# ============================================================
# Slide 5: Architecture
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)
add_rect(slide5, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide5, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide5, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "03  系统架构设计", font_size=32, color=DARK_BLUE, bold=True)

# Architecture layers
layers = [
    ("表示层 (View)", "JSP 页面 + CSS3\n登录/注册、图书列表、借阅管理、后台管理", LIGHT_BLUE),
    ("控制层 (Controller)", "Servlet (10个)\nLoginServlet · RegisterServlet · BorrowServlet · ReturnServlet · IndexServlet · SearchBookServlet · MyBorrowsServlet · AdminBookServlet · AdminBorrowsServlet · LogoutServlet", DARK_BLUE),
    ("业务逻辑层 (Service/DAO)", "UserDao · BookDao · BorrowDao\n封装数据库CRUD操作，处理借阅归还业务逻辑", ACCENT_GREEN),
    ("数据层 (Database)", "MySQL 8.0\nusers表 · books表 · borrow_records表", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(layers):
    y = 2.5 + i * 1.15
    add_rect(slide5, Inches(2.0), Inches(y), Inches(9.3), Inches(1.0), LIGHT_GRAY)
    add_rect(slide5, Inches(2.0), Inches(y), Inches(0.06), Inches(1.0), color)
    add_text_box(slide5, Inches(2.3), Inches(y + 0.1), Inches(3.0), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide5, Inches(5.5), Inches(y + 0.1), Inches(5.5), Inches(0.8),
                 desc, font_size=14, color=DARK_GRAY)

# Arrow between layers
for i in range(3):
    y = 3.5 + i * 1.15
    add_text_box(slide5, Inches(6.3), Inches(y - 0.05), Inches(0.5), Inches(0.3),
                 "▼", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 6: Database Design
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, WHITE)
add_rect(slide6, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide6, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide6, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "04  数据库设计", font_size=32, color=DARK_BLUE, bold=True)

# users table
add_text_box(slide6, Inches(1.3), Inches(2.4), Inches(3.5), Inches(0.4),
             "👤 users（用户表）", font_size=18, color=LIGHT_BLUE, bold=True)
add_multi_text(slide6, Inches(1.3), Inches(2.9), Inches(3.5), Inches(2.5), [
    ("id          INT (PK, AUTO_INCREMENT)", False, DARK_GRAY, 13),
    ("username    VARCHAR(50) UNIQUE", False, DARK_GRAY, 13),
    ("password    VARCHAR(100)", False, DARK_GRAY, 13),
    ("real_name   VARCHAR(50)", False, DARK_GRAY, 13),
    ("phone       VARCHAR(20)", False, DARK_GRAY, 13),
    ("role        ENUM('reader','admin')", False, ACCENT_ORANGE, 13),
], font_size=13)

# books table
add_text_box(slide6, Inches(5.1), Inches(2.4), Inches(4.0), Inches(0.4),
             "📕 books（图书表）", font_size=18, color=ACCENT_GREEN, bold=True)
add_multi_text(slide6, Inches(5.1), Inches(2.9), Inches(4.0), Inches(2.5), [
    ("id                INT (PK)", False, DARK_GRAY, 13),
    ("title             VARCHAR(200)", False, DARK_GRAY, 13),
    ("author            VARCHAR(100)", False, DARK_GRAY, 13),
    ("isbn              VARCHAR(20) UNIQUE", False, DARK_GRAY, 13),
    ("category          VARCHAR(50)", False, DARK_GRAY, 13),
    ("total_copies      INT", False, DARK_GRAY, 13),
    ("available_copies  INT   ← 可借数量", False, ACCENT_ORANGE, 13),
], font_size=13)

# borrow_records table
add_text_box(slide6, Inches(9.5), Inches(2.4), Inches(3.5), Inches(0.4),
             "📋 borrow_records", font_size=18, color=ACCENT_ORANGE, bold=True)
add_multi_text(slide6, Inches(9.5), Inches(2.9), Inches(3.5), Inches(2.5), [
    ("id            INT (PK)", False, DARK_GRAY, 13),
    ("user_id       INT (FK→users)", False, DARK_GRAY, 13),
    ("book_id       INT (FK→books)", False, DARK_GRAY, 13),
    ("borrow_date   DATE", False, DARK_GRAY, 13),
    ("due_date      DATE (借期30天)", False, DARK_GRAY, 13),
    ("return_date   DATE (可空)", False, DARK_GRAY, 13),
    ("status        ENUM('borrowed',\n               'returned','overdue')", False, ACCENT_ORANGE, 13),
], font_size=13)

# ER relationship description
add_rect(slide6, Inches(1.3), Inches(5.6), Inches(11.0), Inches(1.4), LIGHT_GRAY)
add_rect(slide6, Inches(1.3), Inches(5.6), Inches(0.06), Inches(1.4), DARK_BLUE)
add_multi_text(slide6, Inches(1.6), Inches(5.7), Inches(10.5), Inches(1.3), [
    ("🔗 表关系：", True, DARK_BLUE, 16),
    ("users (1) ──→ (N) borrow_records (N) ←── (1) books", False, DARK_GRAY, 15),
    ("借阅记录表通过 user_id 和 book_id 外键关联用户和图书，实现多对多借阅关系。", False, MEDIUM_GRAY, 14),
    ("borrow 时：available_copies - 1  │  return 时：available_copies + 1", False, ACCENT_ORANGE, 14),
])

# ============================================================
# Slide 7: Features Overview
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, WHITE)
add_rect(slide7, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide7, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide7, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "05  功能模块展示", font_size=32, color=DARK_BLUE, bold=True)

# Reader features
add_text_box(slide7, Inches(1.3), Inches(2.4), Inches(5.5), Inches(0.5),
             "👨‍🎓 读者端功能", font_size=22, color=LIGHT_BLUE, bold=True)
reader_feats = [
    "  用户注册与登录",
    "  浏览图书列表（卡片式布局）",
    "  搜索图书（书名 / 作者 / ISBN / 分类）",
    "  一键借阅（自动生成30天借期）",
    "  我的借阅记录（查看 / 还书）",
]
for i, feat in enumerate(reader_feats):
    y = 3.0 + i * 0.5
    add_rect(slide7, Inches(1.3), Inches(y), Inches(5.3), Inches(0.42), LIGHT_GRAY)
    add_text_box(slide7, Inches(1.5), Inches(y + 0.05), Inches(5.0), Inches(0.35),
                 feat, font_size=15, color=DARK_GRAY)

# Admin features
add_text_box(slide7, Inches(7.3), Inches(2.4), Inches(5.5), Inches(0.5),
             "🛠️ 管理员端功能", font_size=22, color=ACCENT_ORANGE, bold=True)
admin_feats = [
    "  图书管理（添加 / 编辑 / 删除）",
    "  模态框表单操作，体验流畅",
    "  查看全部借阅记录",
    "  手动确认归还（超期等特殊情况）",
    "  权限控制：非管理员无法访问后台",
]
for i, feat in enumerate(admin_feats):
    y = 3.0 + i * 0.5
    add_rect(slide7, Inches(7.3), Inches(y), Inches(5.3), Inches(0.42), LIGHT_GRAY)
    add_text_box(slide7, Inches(7.5), Inches(y + 0.05), Inches(5.0), Inches(0.35),
                 feat, font_size=15, color=DARK_GRAY)

# Bottom highlight
add_rect(slide7, Inches(1.3), Inches(5.8), Inches(11.0), Inches(1.1), DARK_BLUE)
add_multi_text(slide7, Inches(1.6), Inches(6.0), Inches(10.4), Inches(0.8), [
    ("🔐 权限控制：", True, WHITE, 16),
    ("所有 Servlet 均校验 Session，普通读者无法访问管理后台，管理员自动跳转至管理页面。", False, RGBColor(0xBB, 0xDE, 0xFB), 14),
])

# ============================================================
# Slide 8: Core Code Showcase
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, WHITE)
add_rect(slide8, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide8, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide8, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "06  核心代码解析", font_size=32, color=DARK_BLUE, bold=True)

# Code box 1 - BorrowServlet
add_text_box(slide8, Inches(1.3), Inches(2.3), Inches(6.0), Inches(0.4),
             "📌 借书核心逻辑（BorrowServlet）", font_size=16, color=LIGHT_BLUE, bold=True)
code_box1 = add_rect(slide8, Inches(1.3), Inches(2.7), Inches(5.8), Inches(2.5), CODE_BG)
add_multi_text(slide8, Inches(1.5), Inches(2.8), Inches(5.4), Inches(2.3), [
    ("// 1. 检查库存", False, RGBColor(0x6A, 0x99, 0x5C), 11),
    ("if (book.getAvailableCopies() <= 0) {", False, RGBColor(0xCC, 0xCC, 0xCC), 11),
    ('    msg = "库存不足，无法借阅！";', False, RGBColor(0xCC, 0xCC, 0xCC), 11),
    ("}", False, RGBColor(0xCC, 0xCC, 0xCC), 11),
    ("// 2. 检查是否已借过该书", False, RGBColor(0x6A, 0x99, 0x5C), 11),
    ("if (borrowDao.hasBorrowed(userId, bookId)) {", False, RGBColor(0xCC, 0xCC, 0xCC), 11),
    ('    msg = "已借过，请勿重复借阅！";', False, RGBColor(0xCC, 0xCC, 0xCC), 11),
    ("}", False, RGBColor(0xCC, 0xCC, 0xCC), 11),
    ("// 3. 创建记录 + 减库存（事务）", False, RGBColor(0x6A, 0x99, 0x5C), 11),
    ("borrowDao.borrow(userId, bookId);", False, RGBColor(0xD4, 0xD4, 0xD4), 11),
    ("bookDao.decreaseAvailable(bookId);", False, RGBColor(0xD4, 0xD4, 0xD4), 11),
], font_size=11)

# Code box 2 - SQL
add_text_box(slide8, Inches(7.5), Inches(2.3), Inches(5.5), Inches(0.4),
             "📌 关键SQL语句", font_size=16, color=ACCENT_GREEN, bold=True)
code_box2 = add_rect(slide8, Inches(7.5), Inches(2.7), Inches(5.5), Inches(2.5), CODE_BG)
add_multi_text(slide8, Inches(7.7), Inches(2.8), Inches(5.1), Inches(2.3), [
    ("-- 借书（自动生成30天借期）", False, RGBColor(0x6A, 0x99, 0x5C), 10),
    ("INSERT INTO borrow_records", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("  (user_id, book_id, borrow_date,", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("   due_date, status)", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("VALUES (?, ?, CURDATE(),", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("  DATE_ADD(CURDATE(),INTERVAL 30 DAY),", False, RGBColor(0xCE, 0x91, 0x78), 10),
    ("  'borrowed')", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("", False, WHITE, 6),
    ("-- 关联查询借阅记录", False, RGBColor(0x6A, 0x99, 0x5C), 10),
    ("SELECT r.*, u.real_name,", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("  b.title, b.author", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("FROM borrow_records r", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("JOIN users u ON r.user_id = u.id", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("JOIN books b ON r.book_id = b.id", False, RGBColor(0x56, 0x9C, 0xD6), 10),
    ("ORDER BY r.status, r.borrow_date", False, RGBColor(0x56, 0x9C, 0xD6), 10),
], font_size=10)

# Architecture highlights
add_text_box(slide8, Inches(1.3), Inches(5.5), Inches(5.8), Inches(0.4),
             "🏗️ 设计亮点", font_size=17, color=DARK_BLUE, bold=True)
add_multi_text(slide8, Inches(1.3), Inches(5.9), Inches(5.8), Inches(1.2), [
    ("• @WebServlet 注解路由，零XML配置", False, DARK_GRAY, 14),
    ("• DAO 模式分离数据访问与业务逻辑", False, DARK_GRAY, 14),
    ("• 数据库连接工具类统一管理资源回收", False, DARK_GRAY, 14),
    ("• Session 会话管理 + 权限拦截", False, DARK_GRAY, 14),
], font_size=14)

add_text_box(slide8, Inches(7.5), Inches(5.5), Inches(5.5), Inches(0.4),
             "🎨 前端亮点", font_size=17, color=DARK_BLUE, bold=True)
add_multi_text(slide8, Inches(7.5), Inches(5.9), Inches(5.5), Inches(1.2), [
    ("• CSS Grid 响应式卡片布局", False, DARK_GRAY, 14),
    ("• 模态框管理表单，页面无跳转", False, DARK_GRAY, 14),
    ("• 渐变色导航栏 + 微动效", False, DARK_GRAY, 14),
    ("• JSTL + EL 表达式动态渲染", False, DARK_GRAY, 14),
], font_size=14)

# ============================================================
# Slide 9: Summary
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, WHITE)
add_rect(slide9, Inches(0), Inches(0), Inches(13.333), Inches(0.08), DARK_BLUE)
add_rect(slide9, Inches(1.0), Inches(1.5), Inches(0.06), Inches(0.6), LIGHT_BLUE)
add_text_box(slide9, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
             "07  项目总结与展望", font_size=32, color=DARK_BLUE, bold=True)

# Summary
add_text_box(slide9, Inches(1.3), Inches(2.5), Inches(5.5), Inches(0.5),
             "✅ 已完成功能", font_size=20, color=ACCENT_GREEN, bold=True)
add_multi_text(slide9, Inches(1.3), Inches(3.1), Inches(5.5), Inches(3.5), [
    ("✓   用户注册 / 登录 / 登出", False, DARK_GRAY, 16),
    ("✓   图书浏览 + 多条件搜索", False, DARK_GRAY, 16),
    ("✓   借阅图书（库存管理 + 重复借阅检测）", False, DARK_GRAY, 16),
    ("✓   归还图书（记录更新 + 库存恢复）", False, DARK_GRAY, 16),
    ("✓   个人借阅历史查看", False, DARK_GRAY, 16),
    ("✓   管理员图书 CRUD 管理", False, DARK_GRAY, 16),
    ("✓   管理员查看全部借阅记录", False, DARK_GRAY, 16),
    ("✓   权限控制 + Session 拦截", False, DARK_GRAY, 16),
    ("✓   嵌入式 Jetty，零配置启动", False, DARK_GRAY, 16),
], font_size=16)

# Future
add_text_box(slide9, Inches(7.3), Inches(2.5), Inches(5.5), Inches(0.5),
             "🔮 后续可扩展方向", font_size=20, color=ACCENT_ORANGE, bold=True)
add_multi_text(slide9, Inches(7.3), Inches(3.1), Inches(5.5), Inches(4.0), [
    ("→  接入 Spring Boot 框架，简化配置", False, DARK_GRAY, 15),
    ("→  增加图书封面图片上传功能", False, DARK_GRAY, 15),
    ("→  引入 MyBatis / JPA 持久层框架", False, DARK_GRAY, 15),
    ("→  增加借阅提醒（短信 / 邮件）", False, DARK_GRAY, 15),
    ("→  添加读者评论与评分系统", False, DARK_GRAY, 15),
    ("→  数据可视化（借阅热力图、统计图表）", False, DARK_GRAY, 15),
    ("→  前后端分离（Vue + REST API）", False, DARK_GRAY, 15),
], font_size=15)

# Bottom stats
add_rect(slide9, Inches(1.3), Inches(6.2), Inches(11.0), Inches(0.9), DARK_BLUE)
stats_items = [
    ("17", "Java源文件"),
    ("10", "Servlet控制器"),
    ("3", "数据表"),
    ("25+", "页面模板文件"),
    ("1", "一键启动命令"),
]
for i, (num, label) in enumerate(stats_items):
    x = 1.5 + i * 2.15
    add_text_box(slide9, Inches(x), Inches(6.3), Inches(1.8), Inches(0.4),
                 num, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, Inches(x), Inches(6.7), Inches(1.8), Inches(0.3),
                 label, font_size=12, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 10: Thank You
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, DARK_BLUE)
add_rect(slide10, Inches(0), Inches(3.5), Inches(13.333), Inches(0.04), LIGHT_BLUE)
add_rect(slide10, Inches(0), Inches(4.3), Inches(13.333), Inches(0.04), RGBColor(0x42, 0xA5, 0xF5))

add_text_box(slide10, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
             "感谢聆听", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide10, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.6),
             "THANK YOU", font_size=28, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
add_text_box(slide10, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6),
             "图书馆借还系统 · JavaWeb期末答辩", font_size=20, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)
add_text_box(slide10, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.6),
             "请各位老师批评指正", font_size=18, color=RGBColor(0x64, 0xB5, 0xF6), alignment=PP_ALIGN.CENTER)

# ===== Save =====
output_path = os.path.join(output_dir, "图书馆借还系统_答辩PPT.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
